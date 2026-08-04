#!/usr/bin/env python3
"""Réécrit la section 5 de l'Objectif 1 en comparaison à protocole identique.

Le tableau opposait deux évaluations qui ne partageaient ni la référence de
boiterie, ni la fenêtre d'analyse, ni la cohorte. Il est remplacé par la
comparaison produite par run_objective1_sls_comparaison_equitable.py : mêmes
14 vaches, même score du 12 mars, même fenêtre de sept jours. L'ancienne lecture
de janvier est conservée sous le tableau, comme trace du diagnostic initial.

Le choix du 12 mars est justifié explicitement, dans le rapport comme dans
l'annexe : les données IceTag commencent après la séance de janvier.

Le script est idempotent.
"""
from __future__ import annotations

import shutil
from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Pt
from pptx import Presentation

PROJET = Path(__file__).resolve().parent
SORTIES = PROJET / "reports" / "objective1_pipeline_icetag"
D1 = Path.home() / "Desktop" / "Livrables_McGill_WellE" / "Objectif1_Pipeline_detection_boiterie"
ANNEXE = D1 / "ANNEXE_pipeline_actuelle_HYPO_instabilite_hybride"

RAPPORT = D1 / "RAPPORTS" / "Objectif1_rapport_livraison.docx"
NOTE = D1 / "NOTES_SOW" / "rapport_validation_concordance.md"
DECK = D1 / "RAPPORTS" / "Objectif1_presentation_detaillee.pptx"
ANNEXE_DOCX = ANNEXE / "RAPPORTS" / "Annexe_pipeline_actuelle_HYPO_instabilite_hybride.docx"

JUSTIFICATION = (
    "Le fichier McGill contient deux séances de notation : une au 15 janvier "
    "2019 et une au 12 mars 2019. Les données IceTag de Winter 2019 débutent le "
    "16 janvier, soit après la première séance : aucune fenêtre capteur ne "
    "précède le score de janvier, alors que 54 jours de données précèdent celui "
    "du 12 mars. La comparaison retient donc le score du 12 mars, seul "
    "compatible avec un protocole de détection."
)

NOTE_JANVIER = (
    "Lecture antérieure conservée pour mémoire : un premier diagnostic comparait "
    "le nombre d'alertes de toute la saison au score du 15 janvier, sur 16 vaches "
    "dont 5 avec SLS >= 2, et ne trouvait aucune concordance (p = 0,649; "
    "rho = 0,033). Cette lecture répond à une question de persistance et non de "
    "détection : trois des cinq vaches boiteuses en janvier ne l'étaient plus au "
    "12 mars. Elle n'est pas comparable au tableau ci-dessus."
)


def _valeurs() -> dict[str, pd.Series]:
    comparaison = pd.read_csv(SORTIES / "objective1_sls_comparaison_equitable.csv")
    return {
        "initiale": comparaison[comparaison.pipeline.str.startswith("Pipeline initiale")].iloc[0],
        "actuelle": comparaison[comparaison.pipeline.str.startswith("Pipeline HYPO")].iloc[0],
    }


def _fr(valeur: float, decimales: int = 3) -> str:
    return f"{valeur:.{decimales}f}".replace(".", ",")


def _paragraphe_apres(document: Document, ancre, texte: str, italique: bool = False):
    """Insère un paragraphe juste après l'élément ancre, au style du corps."""
    modele = next(
        (p for p in document.paragraphs if p.text.strip() and p.style.name == "Normal"),
        None,
    )
    nouveau = document.add_paragraph()
    ancre.addnext(nouveau._element)
    if modele is not None:
        nouveau.style = modele.style
        nouveau.paragraph_format.space_after = modele.paragraph_format.space_after
        nouveau.paragraph_format.alignment = modele.paragraph_format.alignment
    passage = nouveau.add_run(texte)
    passage.font.size = Pt(10)
    passage.font.italic = italique
    return nouveau


def corriger_rapport(v: dict) -> str:
    document = Document(RAPPORT)

    # --- le tableau de la section 5 ---
    cible = None
    for tableau in document.tables:
        entete = [c.text.strip() for c in tableau.rows[0].cells]
        if entete and entete[0] == "Évaluation":
            cible = tableau
            break
    if cible is None:
        return "tableau introuvable"

    lignes = [
        [
            "Évaluation",
            "Cohorte",
            "Résultat",
            "Lecture",
        ],
        [
            "Pipeline initiale IF + règles",
            "14 vaches; 3 avec SLS >= 2",
            f"AUC = {_fr(v['initiale'].auc)}; p = {_fr(v['initiale'].mann_whitney_p)}",
            "Au niveau du hasard sur la fenêtre du score.",
        ],
        [
            "Pipeline HYPO + instabilité + hybride",
            "14 vaches; 3 avec SLS >= 2",
            f"AUC = {_fr(v['actuelle'].auc)}; p = {_fr(v['actuelle'].mann_whitney_p)}",
            "Sépare les deux groupes; strictement exploratoire.",
        ],
    ]
    for index, contenu in enumerate(lignes):
        for colonne, texte in enumerate(contenu):
            cellule = cible.rows[index].cells[colonne]
            paragraphe = cellule.paragraphs[0]
            if paragraphe.runs:
                paragraphe.runs[0].text = texte
                for passage in paragraphe.runs[1:]:
                    passage.text = ""
            else:
                paragraphe.add_run(texte)

    # --- le paragraphe qui suit ---
    ancien_debut = "Des scores SLS synchronisés existent pour Winter 2019."
    nouveau = (
        "Les deux pipelines sont ici évaluées sur un protocole strictement "
        "identique : les mêmes 14 vaches de Winter 2019, la même référence de "
        "boiterie et la même fenêtre de sept jours précédant le score. "
        + JUSTIFICATION
        + " Avec seulement trois vaches SLS >= 2 et un effet du traitement "
        "Exercise confondu avec le statut, aucune sensibilité ni spécificité "
        "robuste ne peut être revendiquée."
    )
    remplace = False
    ancre = None
    for paragraphe in document.paragraphs:
        if paragraphe.text.startswith(ancien_debut) or paragraphe.text.startswith(
            "Les deux pipelines sont ici évaluées"
        ):
            passages = paragraphe.runs
            if not passages:
                continue
            passages[0].text = nouveau
            for passage in passages[1:]:
                passage.text = ""
            ancre = paragraphe._element
            remplace = True
            break
    if ancre is not None and not any(
        p.text.startswith("Lecture antérieure conservée") for p in document.paragraphs
    ):
        _paragraphe_apres(document, ancre, NOTE_JANVIER, italique=True)

    document.core_properties.author = "Aliou Barry"
    document.core_properties.last_modified_by = "Aliou Barry"
    document.save(RAPPORT)
    return "tableau et prose réécrits" if remplace else "prose non trouvée"


def corriger_note(v: dict) -> str:
    texte = NOTE.read_text(encoding="utf-8")
    debut = texte.index("## Concordance exploratoire avec les scores SLS")
    fin = texte.index("## Conclusion")
    bloc = (
        "## Concordance exploratoire avec les scores SLS\n"
        "Les deux pipelines sont comparées sur un protocole strictement identique : "
        "les mêmes 14 vaches de Winter 2019, le score McGill du 12 mars 2019 comme "
        "référence, et les notifications produites dans les sept jours précédant ce "
        "score.\n\n"
        "| Pipeline | Cohorte | AUC | Mann-Whitney p | Spearman rho |\n"
        "|---|---|---:|---:|---:|\n"
        f"| Initiale IF + règles | 14 vaches, 3 avec SLS >= 2 | {_fr(v['initiale'].auc)} | "
        f"{_fr(v['initiale'].mann_whitney_p)} | {_fr(v['initiale'].spearman_rho)} |\n"
        f"| HYPO + instabilité + hybride | 14 vaches, 3 avec SLS >= 2 | {_fr(v['actuelle'].auc)} | "
        f"{_fr(v['actuelle'].mann_whitney_p)} | {_fr(v['actuelle'].spearman_rho)} |\n\n"
        f"{JUSTIFICATION}\n\n"
        f"{NOTE_JANVIER}\n\n"
        "Avec trois cas positifs seulement et un traitement Exercise confondu avec le "
        "statut SLS, ces résultats ne permettent aucune estimation clinique de "
        "sensibilité ou de spécificité.\n\n"
        "Sources : reports/objective1_sls_comparaison_equitable.csv et "
        "objective1_sls_comparaison_cohorte.csv, livrés dans "
        "ANNEXE_pipeline_actuelle_HYPO_instabilite_hybride/TABLEAUX_CSV/.\n\n"
    )
    NOTE.write_text(texte[:debut] + bloc + texte[fin:], encoding="utf-8")
    return "section réécrite"


def corriger_annexe() -> str:
    document = Document(ANNEXE_DOCX)
    if any("deux séances de notation" in p.text for p in document.paragraphs):
        return "déjà justifié"
    ancre = None
    for paragraphe in document.paragraphs:
        if "Le SLS n'a servi ni à entraîner" in paragraphe.text:
            ancre = paragraphe._element
            break
    if ancre is None:
        return "point d'insertion introuvable"
    _paragraphe_apres(document, ancre, JUSTIFICATION)
    document.core_properties.author = "Aliou Barry"
    document.core_properties.last_modified_by = "Aliou Barry"
    document.save(ANNEXE_DOCX)
    return "justification ajoutée"


def corriger_deck(v: dict) -> str:
    presentation = Presentation(DECK)
    diapositive = list(presentation.slides)[9]
    touches = 0
    for forme in diapositive.shapes:
        if not forme.has_text_frame:
            continue
        for paragraphe in forme.text_frame.paragraphs:
            for passage in paragraphe.runs:
                if "16 vaches" in passage.text and "0,649" in passage.text:
                    passage.text = (
                        f"14 vaches; 3 avec SLS ≥ 2 / AUC = {_fr(v['initiale'].auc)}; "
                        f"p = {_fr(v['initiale'].mann_whitney_p)} / Au niveau du hasard."
                    )
                    touches += 1
                elif "14 évaluables" in passage.text:
                    passage.text = (
                        f"14 vaches; 3 avec SLS ≥ 2 / AUC = {_fr(v['actuelle'].auc)}; "
                        f"p = {_fr(v['actuelle'].mann_whitney_p)} / Signal exploratoire."
                    )
                    touches += 1
                elif "Cohortes et protocoles distincts" in passage.text:
                    passage.text = passage.text.replace(
                        "Cohortes et protocoles distincts",
                        "Protocole identique, 3 cas positifs",
                    )
                    touches += 1
                elif passage.text.startswith("Winter 2019 permet deux lectures"):
                    passage.text = (
                        "Comparaison à protocole identique : mêmes vaches, "
                        "même score du 12 mars, même fenêtre de sept jours."
                    )
                    touches += 1

    cadre = diapositive.notes_slide.notes_text_frame
    for paragraphe in cadre.paragraphs:
        for passage in paragraphe.runs:
            if "Pour la pipeline initiale, 16 vaches sont disponibles" in passage.text:
                passage.text = (
                    "Les deux pipelines sont évaluées sur le même protocole : les mêmes "
                    "14 vaches, le score McGill du 12 mars et les notifications des sept "
                    f"jours précédents. La pipeline initiale donne une AUC de "
                    f"{_fr(v['initiale'].auc)} avec p = {_fr(v['initiale'].mann_whitney_p)}, "
                    "soit le niveau du hasard : elle ne produit aucune notification pour "
                    "onze des quatorze vaches sur cette fenêtre."
                )
                touches += 1
            elif "Des scores SLS synchronisés existent bien pour Winter 2019" in passage.text:
                passage.text = "Des scores SLS synchronisés existent pour Winter 2019. " + JUSTIFICATION
                touches += 1
    if touches:
        presentation.save(DECK)
    return f"{touches} passage(s) mis à jour"


def copier_sources() -> str:
    cible = ANNEXE / "TABLEAUX_CSV"
    noms = [
        "objective1_sls_comparaison_equitable.csv",
        "objective1_sls_comparaison_cohorte.csv",
    ]
    for nom in noms:
        shutil.copy2(SORTIES / nom, cible / nom)
    return ", ".join(noms)


def main() -> None:
    valeurs = _valeurs()
    print(f"  sources copiees   : {copier_sources()}")
    print(f"  rapport Word      : {corriger_rapport(valeurs)}")
    print(f"  note de validation: {corriger_note(valeurs)}")
    print(f"  annexe Word       : {corriger_annexe()}")
    print(f"  presentation      : {corriger_deck(valeurs)}")


if __name__ == "__main__":
    main()
