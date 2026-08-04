#!/usr/bin/env python3
"""Met à jour la présentation propre de la réunion du 30 juillet.

Deux corrections :

1. La diapositive 11 annonçait « score SLS du 12 mars » pour les deux approches,
   alors que la ligne IF + règles reposait en réalité sur le score du 15 janvier.
   Elle est remplacée par la comparaison à protocole identique produite par
   run_objective1_sls_comparaison_equitable.py.
2. Douze cadratins subsistaient dans le texte; ils passent en tirets simples ou
   en deux-points, selon leur rôle.

Le script est idempotent.
"""
from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation

PROJET = Path(__file__).resolve().parent
COMPARAISON = (
    PROJET / "reports" / "objective1_pipeline_icetag"
    / "objective1_sls_comparaison_equitable.csv"
)
DECK = (
    Path.home() / "Desktop" / "Livrables_McGill_WellE" / "Reunion_McGill_2026-07-30"
    / "Objectif1_Presentation_McGill_WELL-E.pptx"
)

# Les cadratins servent ici soit de séparateur de plage, soit de définition.
REMPLACEMENTS = {
    " – ": " - ",
    " — ": " : ",
    "—": "-",
    "–": "-",
}


def _fr(valeur: float, decimales: int = 3) -> str:
    return f"{valeur:.{decimales}f}".replace(".", ",")


def corriger_diapo_sls(presentation: Presentation, initiale, moyennes) -> int:
    """Réécrit le bloc IF + règles ligne par ligne, au format du bloc HYPO.

    Le bloc compte quatre paragraphes, une valeur par ligne. Remplacer le seul
    premier passage laisserait les anciennes valeurs sur les lignes suivantes.
    """
    diapositive = list(presentation.slides)[10]
    lignes_if = [
        "14 évaluables; 3 avec SLS ≥ 2",
        f"{_fr(moyennes[0], 2)} vs {_fr(moyennes[1], 2)} notifications",
        f"AUC = {_fr(initiale.auc)}  |  p = {_fr(initiale.mann_whitney_p)}",
        f"ρ = {_fr(initiale.spearman_rho)}  |  p = {_fr(initiale.spearman_p)}",
    ]
    touches = 0
    for forme in diapositive.shapes:
        if not forme.has_text_frame:
            continue
        cadre = forme.text_frame
        entier = cadre.text

        if "16 vaches évaluables" in entier or "AUC = 0,576" in entier:
            for index, paragraphe in enumerate(cadre.paragraphs):
                if index >= len(lignes_if) or not paragraphe.runs:
                    continue
                paragraphe.runs[0].text = lignes_if[index]
                for passage in paragraphe.runs[1:]:
                    passage.text = ""
                touches += 1
            continue

        for paragraphe in cadre.paragraphs:
            for passage in paragraphe.runs:
                texte = passage.text
                if texte.startswith("Winter 2019") and "score SLS du 12 mars" in texte:
                    # Tient sur une seule ligne : au-delà, le texte déborde
                    # sur les deux blocs de résultats.
                    passage.text = (
                        "Winter 2019, score SLS du 12 mars. Même cohorte de "
                        "14 vaches, mêmes sept jours précédant le score."
                    )
                    touches += 1
                elif texte.strip() in ("Aucune association observée", "Au niveau du hasard"):
                    passage.text = "Au niveau du hasard"
                    touches += 1
    return touches


def corriger_cadratins(presentation: Presentation) -> int:
    touches = 0

    def traiter(cadre) -> int:
        n = 0
        for paragraphe in cadre.paragraphs:
            for passage in paragraphe.runs:
                remplace = passage.text
                for ancien, nouveau in REMPLACEMENTS.items():
                    remplace = remplace.replace(ancien, nouveau)
                if remplace != passage.text:
                    passage.text = remplace
                    n += 1
        return n

    for diapositive in presentation.slides:
        for forme in diapositive.shapes:
            if forme.has_text_frame:
                touches += traiter(forme.text_frame)
            if forme.has_table:
                for ligne in forme.table.rows:
                    for cellule in ligne.cells:
                        touches += traiter(cellule.text_frame)
    return touches


def main() -> None:
    comparaison = pd.read_csv(COMPARAISON)
    initiale = comparaison[comparaison.pipeline.str.startswith("Pipeline initiale")].iloc[0]
    moyennes = (initiale.moyenne_sls_ge_2, initiale.moyenne_sls_lt_2)

    presentation = Presentation(DECK)
    print(f"  diapositive 11 : {corriger_diapo_sls(presentation, initiale, moyennes)} passage(s)")
    print(f"  cadratins      : {corriger_cadratins(presentation)} passage(s)")

    proprietes = presentation.core_properties
    proprietes.author = "Aliou Barry"
    proprietes.last_modified_by = "Aliou Barry"
    proprietes.title = "Objectif 1 - Pipeline de détection sur données IceTag"
    presentation.save(DECK)
    print(f"  enregistre     : {DECK.name}")


if __name__ == "__main__":
    main()
