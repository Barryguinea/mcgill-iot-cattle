#!/usr/bin/env python3
"""Corrections de l'audit de l'Objectif 1.

Cinq défauts relevés lors de la revue du paquet de livraison :

1. La présentation livrée portait des traces de l'outil qui l'a produite
   (auteur « Walnut Exporter », thème nommé « ChatGPT », titre « Presentation »).
2. Le tableau de concordance SLS mettait en regard deux évaluations qui
   n'utilisent pas la même référence de boiterie : le score de janvier pour la
   pipeline initiale, celui du 12 mars pour la pipeline actuelle. Quatre vaches
   sur quatorze changent de groupe entre les deux lectures.
3. Les chiffres de la ligne « IF + règles » n'avaient aucune source dans le
   paquet, alors que la ligne HYPO dispose de ses CSV en annexe.
4. La présentation mélangeait les séparateurs décimaux (94.4% et 26,8 %).
5. Le rapport d'annexe portait « python-docx » comme auteur.

Le script est idempotent : il peut être relancé sans dommage.
"""
from __future__ import annotations

import re
import shutil
import zipfile
from pathlib import Path

import pandas as pd
from docx import Document
from pptx import Presentation

PROJET = Path(__file__).resolve().parent
LIVRABLES = Path.home() / "Desktop" / "Livrables_McGill_WellE"
D1 = LIVRABLES / "Objectif1_Pipeline_detection_boiterie"
ANNEXE = D1 / "ANNEXE_pipeline_actuelle_HYPO_instabilite_hybride"
DIAGNOSTIC = PROJET / "reports" / "objective1_pipeline_icetag" / "diagnostic_separabilite"

AUTEUR = "Aliou Barry"
THEME = "McGill WELL-E"


# --------------------------------------------------------------- 1 et 4 : pptx
def nettoyer_presentation(chemin: Path, titre: str) -> dict:
    """Retire les traces de l'outil de génération et uniformise les décimales."""
    bilan = {"metadonnees": 0, "themes": 0, "decimales": 0}

    presentation = Presentation(chemin)
    proprietes = presentation.core_properties
    proprietes.author = AUTEUR
    proprietes.last_modified_by = AUTEUR
    proprietes.title = titre
    proprietes.comments = ""
    proprietes.category = ""
    proprietes.keywords = ""
    bilan["metadonnees"] = 1

    # Décimales : le point devient une virgule dans les nombres, jamais dans les
    # numéros de tâche du SOW (1.1, 1.2) ni dans les noms de fichiers.
    motif = re.compile(r"(?<![\w.])(\d+)\.(\d+)(?=\s*%)")

    def corriger(cadre) -> int:
        touches = 0
        for paragraphe in cadre.paragraphs:
            for passage in paragraphe.runs:
                remplace = motif.sub(r"\1,\2", passage.text)
                if remplace != passage.text:
                    passage.text = remplace
                    touches += 1
        return touches

    for diapositive in presentation.slides:
        for forme in diapositive.shapes:
            if forme.has_text_frame:
                bilan["decimales"] += corriger(forme.text_frame)
            if forme.has_table:
                for ligne in forme.table.rows:
                    for cellule in ligne.cells:
                        bilan["decimales"] += corriger(cellule.text_frame)
    presentation.save(chemin)

    # Le nom du thème vit dans le XML, hors de portée de python-pptx.
    temporaire = chemin.with_suffix(".pptx.tmp")
    with zipfile.ZipFile(chemin) as source:
        with zipfile.ZipFile(temporaire, "w", zipfile.ZIP_DEFLATED) as cible:
            for element in source.infolist():
                contenu = source.read(element.filename)
                if "theme" in element.filename and element.filename.endswith(".xml"):
                    texte = contenu.decode("utf-8")
                    remplace = re.sub(
                        r'(<a:theme[^>]*\sname=")[^"]*(")',
                        rf"\1{THEME}\2",
                        texte,
                    )
                    if remplace != texte:
                        bilan["themes"] += 1
                        contenu = remplace.encode("utf-8")
                elif element.filename == "docProps/app.xml":
                    texte = contenu.decode("utf-8")
                    contenu = re.sub(
                        r"<((?:\w+:)?Application)>[^<]*</\1>",
                        rf"<\1>{AUTEUR}</\1>",
                        texte,
                    ).encode("utf-8")
                cible.writestr(element, contenu)
    temporaire.replace(chemin)
    return bilan


# ------------------------------------------------------------------ 3 : source
def ecrire_source_pipeline_initiale() -> Path:
    """Recopie la cohorte de l'évaluation IF + règles dans l'annexe.

    Sans ce fichier, les chiffres de la première ligne du tableau SLS ne sont
    vérifiables nulle part dans le paquet.
    """
    source = DIAGNOSTIC / "classifieur_loo_predictions.csv"
    cohorte = pd.read_csv(source)
    cohorte = cohorte.rename(
        columns={
            "Cow": "vache",
            "SLS_total": "sls_janvier",
            "boiteuse": "sls_ge_2",
            "proba_predite": "probabilite_predite_loo",
        }
    )
    cohorte["reference_boiterie"] = "score SLS de janvier 2019"
    cible = ANNEXE / "TABLEAUX_CSV" / "pipeline_initiale_validation_sls_cohorte.csv"
    cohorte.to_csv(cible, index=False)

    separabilite = pd.read_csv(DIAGNOSTIC / "separabilite_univariee.csv")
    separabilite.to_csv(
        ANNEXE / "TABLEAUX_CSV" / "pipeline_initiale_separabilite_univariee.csv",
        index=False,
    )
    return cible


# --------------------------------------------------------------- 2 : rédaction
ANCIEN_PARAGRAPHE = (
    "Des scores SLS synchronisés existent pour Winter 2019. Les deux évaluations "
    "utilisent toutefois des cohortes et des protocoles différents; elles ne "
    "constituent pas une comparaison directe de performance."
)

NOUVEAU_PARAGRAPHE = (
    "Des scores SLS synchronisés existent pour Winter 2019. Les deux évaluations "
    "ne reposent cependant pas sur la même référence de boiterie : la première "
    "classe les vaches d'après le score de janvier, la seconde d'après celui du "
    "12 mars. Quatre vaches sur quatorze changent de groupe d'une lecture à "
    "l'autre. Les cohortes et les fenêtres d'analyse diffèrent également. Ces "
    "deux lignes ne constituent donc pas une comparaison directe de performance."
)


def corriger_rapport() -> int:
    chemin = D1 / "RAPPORTS" / "Objectif1_rapport_livraison.docx"
    document = Document(chemin)
    touches = 0
    for paragraphe in document.paragraphs:
        if ANCIEN_PARAGRAPHE.split(";")[0] in paragraphe.text:
            suite = paragraphe.text.split(
                "elles ne constituent pas une comparaison directe de performance."
            )
            reste = suite[1] if len(suite) > 1 else ""
            passages = paragraphe.runs
            if not passages:
                continue
            passages[0].text = NOUVEAU_PARAGRAPHE + reste
            for passage in passages[1:]:
                passage.text = ""
            touches += 1
    if touches:
        document.core_properties.author = AUTEUR
        document.core_properties.last_modified_by = AUTEUR
        document.save(chemin)
    return touches


def corriger_note_validation() -> bool:
    chemin = D1 / "NOTES_SOW" / "rapport_validation_concordance.md"
    texte = chemin.read_text(encoding="utf-8")
    ancien = (
        "Les deux évaluations utilisent des cohortes et des protocoles différents. "
        "Elles ne constituent donc pas une comparaison directe de performance, ni "
        "une estimation clinique de sensibilité ou de spécificité."
    )
    nouveau = (
        "Les deux évaluations ne reposent pas sur la même référence de boiterie : "
        "la première classe les vaches d'après le score SLS de janvier, la seconde "
        "d'après celui du 12 mars. Quatre vaches sur quatorze changent de groupe "
        "d'une lecture à l'autre. Les cohortes et les fenêtres d'analyse diffèrent "
        "également. Ces deux lignes ne constituent donc pas une comparaison directe "
        "de performance, ni une estimation clinique de sensibilité ou de "
        "spécificité.\n\nLa cohorte de la pipeline initiale est livrée dans "
        "ANNEXE_pipeline_actuelle_HYPO_instabilite_hybride/TABLEAUX_CSV/"
        "pipeline_initiale_validation_sls_cohorte.csv; celle de la pipeline "
        "actuelle, dans pipeline_actuelle_validation_sls_cohorte.csv."
    )
    if ancien not in texte:
        return False
    chemin.write_text(texte.replace(ancien, nouveau), encoding="utf-8")
    return True


def corriger_notes_presentation() -> int:
    chemin = D1 / "RAPPORTS" / "Objectif1_presentation_detaillee.pptx"
    presentation = Presentation(chemin)
    ancien = (
        "mais les deux analyses affichées n’utilisent pas exactement la même "
        "cohorte ni le même protocole."
    )
    nouveau = (
        "mais les deux analyses affichées ne reposent pas sur la même référence "
        "de boiterie : la pipeline initiale classe les vaches d’après le score de "
        "janvier, la pipeline actuelle d’après celui du 12 mars. Quatre vaches "
        "sur quatorze changent de groupe d’une lecture à l’autre. Les cohortes et "
        "les fenêtres d’analyse diffèrent également."
    )
    touches = 0
    for diapositive in presentation.slides:
        if not diapositive.has_notes_slide:
            continue
        cadre = diapositive.notes_slide.notes_text_frame
        for paragraphe in cadre.paragraphs:
            for passage in paragraphe.runs:
                if ancien in passage.text:
                    passage.text = passage.text.replace(ancien, nouveau)
                    touches += 1
    if touches:
        presentation.save(chemin)
    return touches


# -------------------------------------------------------------------- 5 : docx
def corriger_auteur_annexe() -> bool:
    chemin = ANNEXE / "RAPPORTS" / "Annexe_pipeline_actuelle_HYPO_instabilite_hybride.docx"
    document = Document(chemin)
    if document.core_properties.author == AUTEUR:
        return False
    document.core_properties.author = AUTEUR
    document.core_properties.last_modified_by = AUTEUR
    document.save(chemin)
    return True


def main() -> None:
    print("1 et 4. Présentation de l'Objectif 1")
    bilan = nettoyer_presentation(
        D1 / "RAPPORTS" / "Objectif1_presentation_detaillee.pptx",
        "Objectif 1 - Pipeline de détection sur données IceTag",
    )
    print(f"   métadonnées réécrites, {bilan['themes']} thèmes renommés, "
          f"{bilan['decimales']} passages de texte passés à la virgule")

    print("3. Source de la ligne « IF + règles »")
    cible = ecrire_source_pipeline_initiale()
    print(f"   écrit : {cible.name} et pipeline_initiale_separabilite_univariee.csv")

    print("2. Formulation du tableau SLS")
    print(f"   rapport Word         : {corriger_rapport()} paragraphe(s)")
    print(f"   note de validation   : {'corrigée' if corriger_note_validation() else 'déjà à jour'}")
    print(f"   notes de la diapo 10 : {corriger_notes_presentation()} passage(s)")

    print("5. Auteur du rapport d'annexe")
    print(f"   {'corrigé' if corriger_auteur_annexe() else 'déjà à jour'}")


if __name__ == "__main__":
    main()
